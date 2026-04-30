import os
import glob
import logging
from pathlib import Path

logging.getLogger("pdfplumber").setLevel(logging.ERROR)
logging.getLogger("pdfminer").setLevel(logging.ERROR)

from src.config import RAW_DIR, SKIP_PATTERNS


def _should_skip(filepath: str) -> bool:
    for pattern in SKIP_PATTERNS:
        if pattern in filepath:
            return True
    return False


def parse_pdf(filepath: str) -> str:
    import pdfplumber
    texts = []
    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
    return "\n\n".join(texts)


def parse_docx(filepath: str) -> str:
    from docx import Document
    doc = Document(filepath)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def parse_pptx(filepath: str) -> str:
    from pptx import Presentation
    prs = Presentation(filepath)
    texts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if len(slide_texts) > 1:
            texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)


def parse_html(filepath: str) -> str:
    from bs4 import BeautifulSoup
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def parse_markdown(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".html": parse_html,
    ".htm": parse_html,
    ".md": parse_markdown,
    ".txt": parse_markdown,
}


def ingest_all() -> dict[str, str]:
    results = {}
    for filepath in sorted(RAW_DIR.rglob("*")):
        if not filepath.is_file():
            continue
        if _should_skip(str(filepath)):
            continue
        ext = filepath.suffix.lower()
        parser = PARSERS.get(ext)
        if parser is None:
            continue
        try:
            text = parser(str(filepath))
            if text.strip():
                results[filepath.name] = text
        except Exception as e:
            print(f"[WARN] 解析失败 {filepath.name}: {e}")
    return results


def list_raw_files() -> list[dict]:
    items = []
    for filepath in sorted(RAW_DIR.rglob("*")):
        if not filepath.is_file() or _should_skip(str(filepath)):
            continue
        ext = filepath.suffix.lower()
        supported = ext in PARSERS
        items.append({
            "name": filepath.name,
            "ext": ext,
            "size_kb": round(filepath.stat().st_size / 1024, 1),
            "supported": supported,
        })
    return items
